import csv
import json
import re
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional


TEXT_EXTENSIONS = {
    ".txt", ".md", ".py", ".js", ".ts", ".html", ".css",
    ".json", ".xml", ".yaml", ".yml", ".sql",
}
SUPPORTED_EXTENSIONS = TEXT_EXTENSIONS | {".csv", ".pdf", ".docx", ".pptx", ".xlsx"}


@dataclass
class Section:
    text: str
    locator: Dict[str, Any]


def _clean_text(text: str) -> str:
    text = text.replace("\x00", "")
    text = re.sub(r"[ \t]+\n", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _parse_text(path: Path) -> List[Section]:
    text = path.read_text("utf-8", errors="ignore")
    if path.suffix.lower() == ".json":
        try:
            text = json.dumps(json.loads(text), ensure_ascii=False, indent=2)
        except json.JSONDecodeError:
            pass
    sections = []
    current_heading: Optional[str] = None
    buffer: List[str] = []
    for line in text.splitlines():
        if re.match(r"^\s{0,3}#{1,6}\s+", line):
            if buffer:
                sections.append(Section(_clean_text("\n".join(buffer)), {"section": current_heading}))
                buffer = []
            current_heading = re.sub(r"^\s{0,3}#{1,6}\s+", "", line).strip()
        buffer.append(line)
    if buffer:
        sections.append(Section(_clean_text("\n".join(buffer)), {"section": current_heading}))
    return [section for section in sections if section.text]


def _parse_csv(path: Path) -> List[Section]:
    rows = []
    with path.open("r", encoding="utf-8-sig", errors="ignore", newline="") as handle:
        reader = csv.reader(handle)
        for row in reader:
            rows.append([str(value) for value in row])
    if not rows:
        return []
    header = rows[0]
    sections = []
    batch_size = 40
    for start in range(1, len(rows), batch_size):
        batch = rows[start:start + batch_size]
        lines = ["\t".join(header), *("\t".join(row) for row in batch)]
        sections.append(Section(
            _clean_text("\n".join(lines)),
            {"sheet": path.stem, "row_start": start + 1, "row_end": start + len(batch)},
        ))
    return sections or [Section("\t".join(header), {"sheet": path.stem, "row_start": 1, "row_end": 1})]


def _parse_pdf(path: Path) -> List[Section]:
    try:
        import fitz
        document = fitz.open(path)
        try:
            return [
                Section(_clean_text(page.get_text("text")), {"page": index + 1})
                for index, page in enumerate(document)
                if _clean_text(page.get_text("text"))
            ]
        finally:
            document.close()
    except ImportError:
        import PyPDF2
        sections = []
        with path.open("rb") as handle:
            for index, page in enumerate(PyPDF2.PdfReader(handle).pages, 1):
                text = _clean_text(page.extract_text() or "")
                if text:
                    sections.append(Section(text, {"page": index}))
        return sections


def _zip_xml_text(archive: zipfile.ZipFile, name: str) -> str:
    from xml.etree import ElementTree
    root = ElementTree.fromstring(archive.read(name))
    return _clean_text("\n".join(
        element.text.strip()
        for element in root.iter()
        if element.tag.endswith("}t") and element.text and element.text.strip()
    ))


def _parse_docx(path: Path) -> List[Section]:
    try:
        from docx import Document
    except ImportError:
        with zipfile.ZipFile(path) as archive:
            return [Section(
                _zip_xml_text(archive, "word/document.xml"),
                {"section": path.stem},
            )]
    document = Document(path)
    sections: List[Section] = []
    heading: Optional[str] = None
    buffer: List[str] = []

    def flush() -> None:
        nonlocal buffer
        text = _clean_text("\n".join(buffer))
        if text:
            sections.append(Section(text, {"section": heading}))
        buffer = []

    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        if paragraph.style and paragraph.style.name.lower().startswith("heading"):
            flush()
            heading = text
        buffer.append(text)
    flush()
    for table_index, table in enumerate(document.tables, 1):
        rows = ["\t".join(cell.text.strip() for cell in row.cells) for row in table.rows]
        text = _clean_text("\n".join(rows))
        if text:
            sections.append(Section(text, {"section": heading, "table": table_index}))
    return sections


def _parse_pptx(path: Path) -> List[Section]:
    try:
        from pptx import Presentation
    except ImportError:
        with zipfile.ZipFile(path) as archive:
            slide_names = sorted(
                (
                    name for name in archive.namelist()
                    if re.fullmatch(r"ppt/slides/slide\d+\.xml", name)
                ),
                key=lambda name: int(re.search(r"(\d+)\.xml$", name).group(1)),
            )
            return [
                Section(_zip_xml_text(archive, name), {"slide": index})
                for index, name in enumerate(slide_names, 1)
                if _zip_xml_text(archive, name)
            ]
    presentation = Presentation(path)
    sections = []
    for index, slide in enumerate(presentation.slides, 1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    texts.append("\t".join(cell.text.strip() for cell in row.cells))
        text = _clean_text("\n".join(texts))
        if text:
            title = slide.shapes.title.text.strip() if slide.shapes.title else None
            sections.append(Section(text, {"slide": index, "section": title}))
    return sections


def _parse_xlsx(path: Path) -> List[Section]:
    try:
        from openpyxl import load_workbook
    except ImportError:
        from xml.etree import ElementTree

        with zipfile.ZipFile(path) as archive:
            shared_strings = []
            if "xl/sharedStrings.xml" in archive.namelist():
                root = ElementTree.fromstring(archive.read("xl/sharedStrings.xml"))
                for item in root:
                    shared_strings.append("".join(
                        node.text or ""
                        for node in item.iter()
                        if node.tag.endswith("}t")
                    ))
            sheet_names = sorted(
                (
                    name for name in archive.namelist()
                    if re.fullmatch(r"xl/worksheets/sheet\d+\.xml", name)
                ),
                key=lambda name: int(re.search(r"(\d+)\.xml$", name).group(1)),
            )
            sections = []
            for index, name in enumerate(sheet_names, 1):
                root = ElementTree.fromstring(archive.read(name))
                rows = []
                for row in root.iter():
                    if not row.tag.endswith("}row"):
                        continue
                    values = []
                    for cell in row:
                        if not cell.tag.endswith("}c"):
                            continue
                        cell_type = cell.attrib.get("t")
                        value_node = next(
                            (node for node in cell if node.tag.endswith("}v")),
                            None,
                        )
                        inline_text = "".join(
                            node.text or ""
                            for node in cell.iter()
                            if node.tag.endswith("}t")
                        )
                        value = inline_text
                        if value_node is not None and value_node.text is not None:
                            value = value_node.text
                            if cell_type == "s":
                                try:
                                    value = shared_strings[int(value)]
                                except (ValueError, IndexError):
                                    pass
                        values.append(value)
                    if values:
                        rows.append(values)
                if not rows:
                    continue
                header = rows[0]
                for start in range(1, len(rows), 40):
                    batch = rows[start:start + 40]
                    text = "\n".join([
                        "\t".join(header),
                        *("\t".join(row) for row in batch),
                    ])
                    sections.append(Section(
                        _clean_text(text),
                        {
                            "sheet": f"Sheet {index}",
                            "row_start": start + 1,
                            "row_end": start + len(batch),
                        },
                    ))
                if len(rows) == 1:
                    sections.append(Section(
                        _clean_text("\t".join(header)),
                        {"sheet": f"Sheet {index}", "row_start": 1, "row_end": 1},
                    ))
            return sections
    workbook = load_workbook(path, read_only=True, data_only=True)
    sections = []
    try:
        for sheet in workbook.worksheets:
            rows = [[str(value) if value is not None else "" for value in row] for row in sheet.iter_rows(values_only=True)]
            rows = [row for row in rows if any(value for value in row)]
            if not rows:
                continue
            header = rows[0]
            for start in range(1, len(rows), 40):
                batch = rows[start:start + 40]
                text = "\n".join(["\t".join(header), *("\t".join(row) for row in batch)])
                sections.append(Section(
                    _clean_text(text),
                    {"sheet": sheet.title, "row_start": start + 1, "row_end": start + len(batch)},
                ))
    finally:
        workbook.close()
    return sections


def parse_document(path: Path) -> List[Section]:
    path = Path(path)
    extension = path.suffix.lower()
    if extension not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"不支持的知识库文件格式: {extension}")
    if extension in TEXT_EXTENSIONS:
        sections = _parse_text(path)
    elif extension == ".csv":
        sections = _parse_csv(path)
    elif extension == ".pdf":
        sections = _parse_pdf(path)
    elif extension == ".docx":
        sections = _parse_docx(path)
    elif extension == ".pptx":
        sections = _parse_pptx(path)
    else:
        sections = _parse_xlsx(path)
    if not sections:
        raise ValueError("文档中没有可索引的文本内容")
    return sections


def _lexical_spans(text: str) -> List[re.Match]:
    return list(re.finditer(r"[\u4e00-\u9fff]|[A-Za-z0-9_]+|[^\s]", text))


def chunk_sections(
    sections: Iterable[Section],
    target_tokens: int = 500,
    overlap_tokens: int = 80,
) -> List[Dict[str, Any]]:
    chunks: List[Dict[str, Any]] = []
    for section in sections:
        spans = _lexical_spans(section.text)
        if not spans:
            continue
        start = 0
        while start < len(spans):
            end = min(start + target_tokens, len(spans))
            content_start = spans[start].start()
            content_end = spans[end - 1].end()
            content = section.text[content_start:content_end].strip()
            if content:
                chunks.append({
                    "content": content,
                    "locator": dict(section.locator),
                    "token_count": end - start,
                })
            if end == len(spans):
                break
            start = max(end - overlap_tokens, start + 1)
    return chunks
