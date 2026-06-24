from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.dml.fill import FillFormat
import re


def parse_markdown(md_content):
    slides = []
    current_slide = {"title": "", "content": []}
    lines = md_content.split('\n')
    
    i = 0
    while i < len(lines):
        line = lines[i]
        
        if line.startswith('## '):
            if current_slide["title"]:
                slides.append(current_slide)
            current_slide = {"title": line[3:], "content": []}
            i += 1
            continue
        
        if line.startswith('### ') or line.startswith('#### '):
            current_slide["content"].append({"type": "subtitle", "text": line.lstrip('# ')})
            i += 1
            continue
        
        if line.startswith('|') and '|' in line:
            cells = [c.strip() for c in line.split('|')]
            current_slide["content"].append({"type": "table_row", "cells": cells})
            i += 1
            continue
        
        if line.startswith('```'):
            code_lang = line[3:].strip() or 'text'
            code_lines = []
            i += 1
            while i < len(lines) and not lines[i].startswith('```'):
                code_lines.append(lines[i])
                i += 1
            i += 1
            current_slide["content"].append({"type": "code", "language": code_lang, "code": '\n'.join(code_lines)})
            continue
        
        if line.startswith('- ') or line.startswith('* '):
            current_slide["content"].append({"type": "bullet", "text": line[2:]})
            i += 1
            continue
        
        if re.match(r'^\d+\.', line):
            match = re.match(r'(\d+)\. (.+)', line)
            if match:
                current_slide["content"].append({"type": "numbered", "num": int(match.group(1)), "text": match.group(2)})
            i += 1
            continue
        
        if line.strip() and not line.startswith('#') and not line.startswith('---'):
            current_slide["content"].append({"type": "paragraph", "text": line.strip()})
        
        i += 1
    
    if current_slide["title"]:
        slides.append(current_slide)
    
    return slides


def set_gradient_background(slide, colors):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors[0]


def add_decoration_shapes(slide):
    circle1 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        left=Inches(9), top=Inches(-1),
        width=Inches(4), height=Inches(4)
    )
    circle1.fill.solid()
    circle1.fill.fore_color.rgb = RGBColor(59, 130, 246)
    circle1.line.fill.background()
    circle1.shadow.inherit = False
    
    circle2 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        left=Inches(-1), top=Inches(4),
        width=Inches(3), height=Inches(3)
    )
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = RGBColor(168, 85, 247)
    circle2.line.fill.background()
    circle2.shadow.inherit = False
    
    rec1 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=Inches(0), top=Inches(6.8),
        width=Inches(13.33), height=Inches(0.7)
    )
    rec1.fill.solid()
    rec1.fill.fore_color.rgb = RGBColor(59, 130, 246)
    rec1.line.fill.background()


def add_title_slide(prs, title, subtitle=""):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    set_gradient_background(slide, [RGBColor(248, 250, 252)])
    
    circle_bg = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        left=Inches(5), top=Inches(-3),
        width=Inches(10), height=Inches(10)
    )
    circle_bg.fill.solid()
    circle_bg.fill.fore_color.rgb = RGBColor(239, 246, 255)
    circle_bg.line.fill.background()
    
    title_shape = slide.shapes.title
    title_shape.text = ""
    title_textbox = slide.shapes.add_textbox(
        left=Inches(1), top=Inches(2),
        width=Inches(11.33), height=Inches(2)
    )
    tf = title_textbox.text_frame
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER
    
    subtitle_textbox = slide.shapes.add_textbox(
        left=Inches(2), top=Inches(4),
        width=Inches(9.33), height=Inches(1)
    )
    tf = subtitle_textbox.text_frame
    p = tf.add_paragraph()
    p.text = subtitle
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(102, 102, 102)
    p.alignment = PP_ALIGN.CENTER
    
    logo_box = slide.shapes.add_textbox(
        left=Inches(11), top=Inches(6.5),
        width=Inches(2), height=Inches(0.5)
    )
    p = logo_box.text_frame.add_paragraph()
    p.text = "AI Assistant"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(59, 130, 246)
    p.font.bold = True
    
    return slide


def add_content_slide(prs, title, content_items):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    set_gradient_background(slide, [RGBColor(248, 250, 252)])
    add_decoration_shapes(slide)
    
    title_box = slide.shapes.add_textbox(
        left=Inches(0.8), top=Inches(0.5),
        width=Inches(12), height=Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    
    title_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=Inches(0.8), top=Inches(1.4),
        width=Inches(3), height=Inches(0.08)
    )
    title_line.fill.solid()
    title_line.fill.fore_color.rgb = RGBColor(59, 130, 246)
    title_line.line.fill.background()
    
    content_textbox = slide.shapes.add_textbox(
        left=Inches(0.8), top=Inches(1.8),
        width=Inches(11.5), height=Inches(4.5)
    )
    tf = content_textbox.text_frame
    tf.word_wrap = True
    
    for item in content_items:
        if item["type"] == "subtitle":
            p = tf.add_paragraph()
            p.text = item["text"]
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = RGBColor(59, 130, 246)
            p.space_before = Pt(15)
        elif item["type"] == "bullet":
            p = tf.add_paragraph()
            p.text = "• " + item["text"]
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(51, 51, 51)
            p.level = 0
            p.space_before = Pt(6)
        elif item["type"] == "numbered":
            p = tf.add_paragraph()
            p.text = f"{item['num']}. {item['text']}"
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(51, 51, 51)
            p.level = 0
            p.space_before = Pt(6)
        elif item["type"] == "paragraph":
            p = tf.add_paragraph()
            p.text = item["text"]
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(68, 68, 68)
            p.alignment = PP_ALIGN.LEFT
            p.space_before = Pt(8)
    
    return slide


def add_table_slide(prs, title, rows):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    set_gradient_background(slide, [RGBColor(248, 250, 252)])
    add_decoration_shapes(slide)
    
    title_box = slide.shapes.add_textbox(
        left=Inches(0.8), top=Inches(0.5),
        width=Inches(12), height=Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    
    title_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=Inches(0.8), top=Inches(1.4),
        width=Inches(3), height=Inches(0.08)
    )
    title_line.fill.solid()
    title_line.fill.fore_color.rgb = RGBColor(59, 130, 246)
    title_line.line.fill.background()
    
    if len(rows) < 2:
        return slide
    
    num_rows = len(rows)
    num_cols = len(rows[0])
    
    left = Inches(1)
    top = Inches(2)
    width = Inches(11.33)
    height = Inches(0.65 * num_rows)
    
    table = slide.shapes.add_table(num_rows, num_cols, left, top, width, height).table
    table.rows[0].height = Inches(0.7)
    
    for i, row in enumerate(rows):
        for j, cell in enumerate(row):
            table.cell(i, j).text = cell
            paragraph = table.cell(i, j).text_frame.paragraphs[0]
            paragraph.font.size = Pt(12)
            paragraph.alignment = PP_ALIGN.CENTER
            
            if i == 0:
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(255, 255, 255)
                table.cell(i, j).fill.solid()
                table.cell(i, j).fill.fore_color.rgb = RGBColor(0, 51, 102)
            else:
                paragraph.font.color.rgb = RGBColor(51, 51, 51)
                table.cell(i, j).fill.solid()
                table.cell(i, j).fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    for j in range(num_cols):
        table.columns[j].width = Inches(width / num_cols)
    
    return slide


def add_code_slide(prs, title, language, code):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    set_gradient_background(slide, [RGBColor(248, 250, 252)])
    add_decoration_shapes(slide)
    
    title_box = slide.shapes.add_textbox(
        left=Inches(0.8), top=Inches(0.5),
        width=Inches(12), height=Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    
    title_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=Inches(0.8), top=Inches(1.4),
        width=Inches(3), height=Inches(0.08)
    )
    title_line.fill.solid()
    title_line.fill.fore_color.rgb = RGBColor(59, 130, 246)
    title_line.line.fill.background()
    
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left=Inches(0.8), top=Inches(1.8),
        width=Inches(11.7), height=Inches(4.5)
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = RGBColor(26, 26, 26)
    bg_shape.line.fill.background()
    
    lang_tag = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=Inches(0.8), top=Inches(1.8),
        width=Inches(2), height=Inches(0.4)
    )
    lang_tag.fill.solid()
    lang_tag.fill.fore_color.rgb = RGBColor(59, 130, 246)
    lang_tag.line.fill.background()
    
    lang_textbox = slide.shapes.add_textbox(
        left=Inches(1), top=Inches(1.85),
        width=Inches(1.5), height=Inches(0.3)
    )
    p = lang_textbox.text_frame.add_paragraph()
    p.text = language.upper()
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    code_textbox = slide.shapes.add_textbox(
        left=Inches(1), top=Inches(2.3),
        width=Inches(11.3), height=Inches(4)
    )
    tf = code_textbox.text_frame
    tf.word_wrap = True
    
    code_lines = code.split('\n')[:30]
    for idx, line in enumerate(code_lines):
        p = tf.add_paragraph()
        p.text = line if line else ' '
        p.font.size = Pt(9)
        p.font.name = "Consolas"
        
        line_stripped = line.strip()
        
        if line_stripped.startswith('#') or line_stripped.startswith('//') or line_stripped.startswith('/*') or line_stripped.startswith('*'):
            p.font.color.rgb = RGBColor(107, 114, 128)
            p.font.italic = True
        elif line_stripped.startswith('"""') or line_stripped.startswith("'''"):
            p.font.color.rgb = RGBColor(107, 114, 128)
            p.font.italic = True
        elif 'def ' in line or 'class ' in line or 'function ' in line or 'async ' in line or 'const ' in line or 'let ' in line or 'var ' in line:
            p.font.color.rgb = RGBColor(168, 85, 247)
            p.font.bold = True
        elif '"' in line or "'" in line:
            p.font.color.rgb = RGBColor(34, 197, 94)
        elif 'import ' in line or 'from ' in line or 'require(' in line or 'export ' in line:
            p.font.color.rgb = RGBColor(139, 92, 246)
        elif any(kw in line for kw in ['return ', 'if ', 'else ', 'for ', 'while ', 'await ', 'async ', 'new ']):
            p.font.color.rgb = RGBColor(59, 130, 246)
            p.font.bold = True
        elif line_stripped.startswith('@'):
            p.font.color.rgb = RGBColor(245, 158, 11)
        else:
            p.font.color.rgb = RGBColor(226, 232, 240)
        
        p.space_before = Pt(0)
    
    return slide


def add_architecture_slide(prs, title):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    set_gradient_background(slide, [RGBColor(248, 250, 252)])
    add_decoration_shapes(slide)
    
    title_box = slide.shapes.add_textbox(
        left=Inches(0.8), top=Inches(0.5),
        width=Inches(12), height=Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    
    title_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=Inches(0.8), top=Inches(1.4),
        width=Inches(3), height=Inches(0.08)
    )
    title_line.fill.solid()
    title_line.fill.fore_color.rgb = RGBColor(59, 130, 246)
    title_line.line.fill.background()
    
    layers = [
        ("前端应用", ["Vue 3", "Vite", "Pinia"], RGBColor(59, 130, 246)),
        ("后端 API", ["FastAPI", "Python", "LangChain"], RGBColor(34, 197, 94)),
        ("数据层", ["MySQL", "ChromaDB", "MinIO"], RGBColor(245, 158, 11)),
        ("AI 模型", ["DeepSeek", "BGE-M3", "Vision"], RGBColor(168, 85, 247)),
    ]
    
    layer_height = Inches(1.2)
    start_top = Inches(2)
    
    for i, (layer_name, items, color) in enumerate(layers):
        top = start_top + i * (layer_height + Inches(0.3))
        
        layer_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left=Inches(1), top=top,
            width=Inches(11.33), height=layer_height
        )
        layer_bg.fill.solid()
        layer_bg.fill.fore_color.rgb = color
        layer_bg.line.color.rgb = RGBColor(255, 255, 255)
        layer_bg.line.width = Pt(2)
        
        icon_box = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            left=Inches(1.5), top=top + Inches(0.3),
            width=Inches(0.6), height=Inches(0.6)
        )
        icon_box.fill.solid()
        icon_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
        icon_box.line.fill.background()
        
        name_textbox = slide.shapes.add_textbox(
            left=Inches(2.5), top=top + Inches(0.25),
            width=Inches(4), height=Inches(0.7)
        )
        name_p = name_textbox.text_frame.add_paragraph()
        name_p.text = layer_name
        name_p.font.size = Pt(18)
        name_p.font.bold = True
        name_p.font.color.rgb = RGBColor(255, 255, 255)
        
        items_text = " | ".join(items)
        items_textbox = slide.shapes.add_textbox(
            left=Inches(7), top=top + Inches(0.35),
            width=Inches(5), height=Inches(0.5)
        )
        items_p = items_textbox.text_frame.add_paragraph()
        items_p.text = items_text
        items_p.font.size = Pt(12)
        items_p.font.color.rgb = RGBColor(255, 255, 255)
        
        if i < len(layers) - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.DOWN_ARROW,
                left=Inches(6.3), top=top + layer_height + Inches(0.05),
                width=Inches(0.7), height=Inches(0.25)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(156, 163, 175)
            arrow.line.fill.background()
    
    return slide


def add_directory_slide(prs, title):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    set_gradient_background(slide, [RGBColor(248, 250, 252)])
    add_decoration_shapes(slide)
    
    title_box = slide.shapes.add_textbox(
        left=Inches(0.8), top=Inches(0.5),
        width=Inches(12), height=Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    
    title_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=Inches(0.8), top=Inches(1.4),
        width=Inches(3), height=Inches(0.08)
    )
    title_line.fill.solid()
    title_line.fill.fore_color.rgb = RGBColor(59, 130, 246)
    title_line.line.fill.background()
    
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left=Inches(0.8), top=Inches(1.8),
        width=Inches(11.7), height=Inches(4.5)
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = RGBColor(26, 26, 26)
    bg_shape.line.fill.background()
    
    textbox = slide.shapes.add_textbox(
        left=Inches(1), top=Inches(2),
        width=Inches(11.3), height=Inches(4.1)
    )
    tf = textbox.text_frame
    tf.word_wrap = True
    
    directory = [
        "demo1/",
        "├── app.py                    # 后端主入口",
        "├── config.py                 # 配置管理",
        "├── redis_client.py           # Redis客户端",
        "├── kb_worker.py              # 任务队列工作器",
        "├── frontend/",
        "│   ├── src/views/            # 页面组件",
        "│   ├── src/api/              # API调用",
        "│   └── src/stores/           # 状态管理",
        "├── rag/",
        "│   ├── embeddings.py         # 嵌入模型",
        "│   ├── knowledge_base.py     # 知识库管理",
        "│   ├── retriever.py          # 检索器",
        "│   └── generator.py          # 生成器",
        "└── rag/data/                 # 数据目录",
    ]
    
    for line in directory:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(11)
        p.font.name = "Consolas"
        
        if ".py" in line:
            p.font.color.rgb = RGBColor(96, 165, 250)
        elif ".vue" in line:
            p.font.color.rgb = RGBColor(139, 92, 246)
        elif "/" in line and not ".py" in line and not ".vue" in line:
            p.font.color.rgb = RGBColor(226, 232, 240)
        elif "#" in line:
            p.font.color.rgb = RGBColor(107, 114, 128)
    
    return slide


def add_er_diagram_slide(prs, title):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    set_gradient_background(slide, [RGBColor(248, 250, 252)])
    add_decoration_shapes(slide)
    
    title_box = slide.shapes.add_textbox(
        left=Inches(0.8), top=Inches(0.5),
        width=Inches(12), height=Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    
    title_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=Inches(0.8), top=Inches(1.4),
        width=Inches(3), height=Inches(0.08)
    )
    title_line.fill.solid()
    title_line.fill.fore_color.rgb = RGBColor(59, 130, 246)
    title_line.line.fill.background()
    
    entities = [
        ("users", "用户表", 0.5, 2),
        ("conversations", "会话表", 4, 2),
        ("messages", "消息表", 8, 2),
        ("knowledge_collections", "知识库集合", 0.5, 4.5),
        ("documents", "文档表", 4, 4.5),
        ("chunks\n(in ChromaDB)", "文档块", 8, 4.5),
    ]
    
    for name, desc, x, y in entities:
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left=Inches(x), top=Inches(y),
            width=Inches(2.5), height=Inches(1)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(59, 130, 246)
        box.line.fill.background()
        
        tb = slide.shapes.add_textbox(
            left=Inches(x), top=Inches(y + 0.15),
            width=Inches(2.5), height=Inches(0.7)
        )
        tf = tb.text_frame
        p = tf.add_paragraph()
        p.text = name
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(9)
        p2.font.color.rgb = RGBColor(226, 232, 240)
        p2.alignment = PP_ALIGN.CENTER
    
    connections = [
        (3, 2.5, 4, 2.5),
        (6.5, 2.5, 8, 2.5),
        (3, 5, 4, 5),
        (6.5, 5, 8, 5),
        (1.75, 3, 1.75, 4.5),
        (5.25, 3, 5.25, 4.5),
    ]
    
    for x1, y1, x2, y2 in connections:
        line = slide.shapes.add_connector(
            1, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
        )
        line.line.color.rgb = RGBColor(156, 163, 175)
        line.line.width = Pt(1.5)
    
    return slide


def add_flow_slide(prs, title, content):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    set_gradient_background(slide, [RGBColor(248, 250, 252)])
    add_decoration_shapes(slide)
    
    title_box = slide.shapes.add_textbox(
        left=Inches(0.8), top=Inches(0.5),
        width=Inches(12), height=Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    
    title_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=Inches(0.8), top=Inches(1.4),
        width=Inches(3), height=Inches(0.08)
    )
    title_line.fill.solid()
    title_line.fill.fore_color.rgb = RGBColor(59, 130, 246)
    title_line.line.fill.background()
    
    steps = []
    for item in content:
        if item["type"] == "paragraph":
            text = item["text"].strip()
            if "──→" in text or "→" in text or "│" in text or "▼" in text or "└─" in text or "├─" in text or text.startswith("1.") or text.startswith("2.") or text.startswith("3.") or text.startswith("4.") or text.startswith("5.") or text.startswith("6."):
                steps.append(text)
        elif item["type"] == "bullet":
            steps.append(item["text"])
    
    if not steps:
        for item in content:
            if item["type"] == "paragraph":
                steps.append(item["text"])
    
    step_height = Inches(0.7)
    start_top = Inches(2)
    
    for i, step in enumerate(steps[:6]):
        top = start_top + i * (step_height + Inches(0.2))
        
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            left=Inches(0.8), top=top,
            width=Inches(0.6), height=Inches(0.6)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = RGBColor(59, 130, 246)
        circle.line.fill.background()
        
        num_tb = slide.shapes.add_textbox(
            left=Inches(0.8), top=top + Inches(0.1),
            width=Inches(0.6), height=Inches(0.4)
        )
        p = num_tb.text_frame.add_paragraph()
        p.text = str(i + 1)
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        step_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left=Inches(1.6), top=top,
            width=Inches(10.5), height=step_height
        )
        step_box.fill.solid()
        step_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
        step_box.line.color.rgb = RGBColor(59, 130, 246)
        step_box.line.width = Pt(1)
        
        step_tb = slide.shapes.add_textbox(
            left=Inches(1.8), top=top + Inches(0.15),
            width=Inches(10.2), height=step_height - Inches(0.1)
        )
        tf = step_tb.text_frame
        p = tf.add_paragraph()
        display_text = step.replace("──→", "→").replace("│", "").replace("▼", "↓").replace("└─", "").replace("├─", "").replace("├──", "").replace("└──", "")
        p.text = display_text.strip()
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.alignment = PP_ALIGN.LEFT
    
    return slide


def md_to_ppt(md_file, ppt_file):
    with open(md_file, 'r', encoding='utf-8') as f:
        md_content = f.read()
    
    slides_data = parse_markdown(md_content)
    
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    add_title_slide(prs, "AI 智能助手项目介绍", "基于大语言模型的智能对话系统")
    
    for slide_data in slides_data:
        title = slide_data["title"]
        content = slide_data["content"]
        
        has_table = any(item["type"] == "table_row" for item in content)
        has_code = any(item["type"] == "code" for item in content)
        
        if "技术架构" in title or "架构图" in title or "整体架构" in title:
            add_architecture_slide(prs, title)
        elif "文件结构" in title or "目录" in title or "目录结构" in title:
            add_directory_slide(prs, title)
        elif "实体关系" in title:
            add_er_diagram_slide(prs, title)
        elif "工作流程" in title or "上传流程" in title or "启动顺序" in title:
            add_flow_slide(prs, title, content)
        elif has_table and not has_code:
            rows = [item["cells"] for item in content if item["type"] == "table_row"]
            add_table_slide(prs, title, rows)
        elif has_code:
            code_items = [item for item in content if item["type"] == "code"]
            for idx, code_item in enumerate(code_items):
                if idx == 0:
                    code_title = title
                else:
                    code_title = f"{title} ({idx+1})"
                add_code_slide(prs, code_title, code_item["language"], code_item["code"])
        else:
            add_content_slide(prs, title, content)
    
    prs.save(ppt_file)
    print(f"PPT文件已保存到: {ppt_file}")


if __name__ == "__main__":
    md_file = "项目介绍PPT.md"
    ppt_file = "项目介绍PPT.pptx"
    md_to_ppt(md_file, ppt_file)
